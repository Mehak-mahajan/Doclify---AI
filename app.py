
import streamlit as st
import pdfplumber
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from groq import Groq
from dotenv import load_dotenv
import os
import docx
import pptx
from langdetect import detect

load_dotenv()

# ---- PAGE SETUP ----
st.set_page_config(
    page_title="Doclify AI",
    page_icon="📚",
    layout="centered"
)

st.title("📚 Doclify AI")
st.subheader("Upload your notes and chat with them!")

# ---- LOAD MODELS ----
@st.cache_resource
def load_model():
    return SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")

embedding_model = load_model()
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# ---- TEXT EXTRACTION WITH PAGE NUMBERS ----
def extract_text(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()
    pages = []

    if file_type == "pdf":
        with pdfplumber.open(uploaded_file) as pdf:
            total_pages = len(pdf.pages)
            progress = st.progress(0)
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    pages.append((page_text, i + 1))
                progress.progress((i + 1) / total_pages)
            progress.empty()

    elif file_type == "txt":
        text = uploaded_file.read().decode("utf-8")
        pages.append((text, 1))

    elif file_type == "docx":
        doc = docx.Document(uploaded_file)
        full_text = ""
        for paragraph in doc.paragraphs:
            full_text += paragraph.text + "\n"
        pages.append((full_text, 1))

    elif file_type == "pptx":
        prs = pptx.Presentation(uploaded_file)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text + "\n"
            if slide_text:
                pages.append((slide_text, i + 1))

    elif file_type == "csv":
        text = uploaded_file.read().decode("utf-8")
        pages.append((text, 1))

    return pages

# ---- PROCESS FILE ----
def process_file(uploaded_file):
    # File size warning
    file_size = uploaded_file.size / (1024 * 1024)
    if file_size > 20:
        st.warning(f"⚠️ Large file ({file_size:.1f}MB) — processing may take longer!")

    pages = extract_text(uploaded_file)

    if not pages:
        return None, None, None

    chunks = []
    metadata = []

    # Larger chunk for large files
    chunk_size = 1000 if file_size > 10 else 500

    for page_text, page_num in pages:
        for i in range(0, len(page_text), chunk_size):
            chunk = page_text[i:i+chunk_size]
            chunks.append(chunk)
            metadata.append({
                "file": uploaded_file.name,
                "page": page_num
            })

    if not chunks:
        return None, None, None

    vectors = embedding_model.encode(
        chunks,
        show_progress_bar=False,
        batch_size=32
    )

    dimension = len(vectors[0])
    index = faiss.IndexFlatL2(dimension)
    index.add(np.array(vectors))

    return chunks, metadata, index

# ---- SUMMARIZE DOCUMENT ----
def summarize_document(chunks):
    # Take first 20 chunks for summary
    # to avoid token limit
    sample_chunks = chunks[:20]
    combined_text = "\n\n".join(sample_chunks)

    response = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "system",
                "content": "You are a helpful study assistant. Summarize the provided document clearly and concisely."
            },
            {
                "role": "user",
                "content": f"Please summarize this document in bullet points covering main topics:\n\n{combined_text}"
            }
        ],
        stream=True
    )

    for chunk in response:
        content = chunk.choices[0].delta.content
        if content:
            yield content

# ---- ASK QUESTION WITH STREAMING + CITATIONS ----
def ask_question(question, chunks, metadata, index):
    # Detect language
    try:
        lang = detect(question)
    except:
        lang = "en"

    question_vector = embedding_model.encode([question])
    distances, indices = index.search(np.array(question_vector), 3)

    context = ""
    sources = []

    for i in indices[0]:
        context += chunks[i] + "\n\n"
        source = metadata[i]
        if source not in sources:
            sources.append(source)

    stream = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "system",
                "content": f"You are a helpful study assistant. Answer in the same language as the question (detected: {lang}). Use ONLY the context provided. If answer not found say 'I could not find this in the provided notes.'"
            },
            {
                "role": "user",
                "content": f"Context:\n{context}\n\nQuestion: {question}"
            }
        ],
        stream=True
    )

    for chunk in stream:
        content = chunk.choices[0].delta.content
        if content:
            yield content, sources

# ---- MAIN UI ----
uploaded_files = st.file_uploader(
    "Upload your notes (PDF, Word, TXT, PPT, CSV)",
    type=["pdf", "txt", "docx", "pptx", "csv"],
    accept_multiple_files=True
)

if uploaded_files:
    with st.spinner("Processing your documents..."):
        all_chunks = []
        all_metadata = []

        for uploaded_file in uploaded_files:
            result = process_file(uploaded_file)

            if result[0] is not None:
                chunks, metadata, index = result
                all_chunks.extend(chunks)
                all_metadata.extend(metadata)
                st.write(f"✅ {uploaded_file.name} → {len(chunks)} chunks")
            else:
                st.warning(f"⚠️ Could not extract text from {uploaded_file.name}")

        if all_chunks:
            vectors = embedding_model.encode(
                all_chunks,
                show_progress_bar=False,
                batch_size=32
            )
            dimension = len(vectors[0])
            combined_index = faiss.IndexFlatL2(dimension)
            combined_index.add(np.array(vectors))

    st.success(f"✅ {len(uploaded_files)} file(s) processed! "
               f"Total {len(all_chunks)} chunks ready!")

    # ---- SUMMARIZE BUTTON ----
    if st.button("📝 Summarize Document"):
        with st.expander("📝 Document Summary", expanded=True):
            st.write_stream(summarize_document(all_chunks))

    st.divider()

    # ---- CHAT SECTION ----
    if "messages" not in st.session_state:
        st.session_state.messages = []

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])
            if "sources" in message:
                for source in message["sources"]:
                    st.caption(
                        f"📄 Source: {source['file']} | "
                        f"Page {source['page']}"
                    )

    question = st.chat_input("Ask anything about your notes...")

    if question:
        with st.chat_message("user"):
            st.write(question)

        st.session_state.messages.append({
            "role": "user",
            "content": question
        })

        with st.chat_message("assistant"):
            answer_text = ""
            answer_sources = []
            answer_placeholder = st.empty()

            for content, sources in ask_question(
                question, all_chunks, all_metadata, combined_index
            ):
                answer_text += content
                answer_sources = sources
                answer_placeholder.write(answer_text)

            for source in answer_sources:
                st.caption(
                    f"📄 Source: {source['file']} | "
                    f"Page {source['page']}"
                )

        st.session_state.messages.append({
            "role": "assistant",
            "content": answer_text,
            "sources": answer_sources
        })

else:
    st.info("👆 Upload PDF, Word, TXT, PowerPoint or CSV files to get started!")